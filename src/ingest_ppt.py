import os
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import spacy
import nltk
from nltk.corpus import stopwords
from sentence_transformers import SentenceTransformer
import chromadb
from tqdm import tqdm
from loguru import logger
from PIL import Image
import io

import sys
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
from src.suppress_warnings import *
sys.path.insert(0, str(Path(__file__).parent.parent))
from config import (
    PPT_DIR, IMAGE_DIR, EMBEDDINGS_DIR, VECTOR_DB_DIR,
    TEXT_EMBED_MODEL, CHROMA_COLLECTION
)

_nlp      = None
_embedder = None
_stopwords = None

def get_nlp():
    global _nlp
    if _nlp is None:
        try:
            _nlp = spacy.load("en_core_web_sm")
        except OSError:
            logger.warning("Run: python -m spacy download en_core_web_sm")
    return _nlp

def get_embedder():
    global _embedder
    if _embedder is None:
        logger.info(f"Loading embedder: {TEXT_EMBED_MODEL}")
        _embedder = SentenceTransformer(TEXT_EMBED_MODEL)
    return _embedder

def get_stopwords():
    global _stopwords
    if _stopwords is None:
        try:
            _stopwords = set(stopwords.words("english"))
        except LookupError:
            nltk.download("stopwords")
            _stopwords = set(stopwords.words("english"))
    return _stopwords

def extract_slide_texts(pptx_path: Path) -> list[dict]:
    prs    = Presentation(str(pptx_path))
    slides = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = ""
        body  = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title:
                title = text
            else:
                body.append(text)
        full_text = (title + " " + " ".join(body)).strip()
        if full_text:
            slides.append({
                "slide_num" : slide_num,
                "title"     : title,
                "body"      : " ".join(body),
                "full_text" : full_text,
            })
    return slides

def extract_topics(text: str) -> list[str]:
    nlp = get_nlp()
    if nlp is None:
        return []
    stop   = get_stopwords()
    doc    = nlp(text[:20_000])
    topics = set()
    for ent in doc.ents:
        if ent.label_ in {"ORG", "PERSON", "GPE", "PRODUCT", "EVENT"}:
            topics.add(ent.text.strip().lower())
    for chunk in doc.noun_chunks:
        phrase = chunk.text.strip().lower()
        words  = [w for w in phrase.split() if w not in stop and len(w) > 2]
        if len(words) >= 1 and len(phrase) > 3:
            topics.add(phrase)
    return list(topics)[:30]

def extract_images_from_pptx(pptx_path: Path, doc_name: str) -> list[str]:
    prs   = Presentation(str(pptx_path))
    saved = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:
                try:
                    image      = shape.image
                    img_bytes  = image.blob
                    ext        = image.ext
                    if len(img_bytes) < 5000:
                        continue
                    fname = f"{doc_name}_slide{slide_num}_img{shape_idx}.{ext}"
                    fpath = IMAGE_DIR / fname
                    with open(fpath, "wb") as f:
                        f.write(img_bytes)
                    saved.append(str(fpath))
                except Exception as e:
                    logger.warning(f"Could not extract image from {doc_name} slide {slide_num}: {e}")

    return saved

def store_in_chromadb(doc_name: str, pptx_path: str,
                       slides: list[dict], embeddings: list,
                       topics_per_slide: list):
    client     = chromadb.PersistentClient(path=str(VECTOR_DB_DIR))
    collection = client.get_or_create_collection(
                     name=CHROMA_COLLECTION,
                     metadata={"hnsw:space": "cosine"}
                 )
    ids, docs, metas, embs = [], [], [], []
    for slide, emb, topics in zip(slides, embeddings, topics_per_slide):
        ids.append(f"ppt_{doc_name}_slide{slide['slide_num']}")
        docs.append(slide["full_text"])
        embs.append(emb)
        metas.append({
            "source"      : str(pptx_path),
            "source_type" : "ppt",
            "doc_name"    : doc_name,
            "slide_num"   : slide["slide_num"],
            "slide_title" : slide["title"],
            "topics"      : json.dumps(topics),
        })
    for b in range(0, len(ids), 100):
        collection.upsert(
            ids        = ids[b:b+100],
            embeddings = embs[b:b+100],
            documents  = docs[b:b+100],
            metadatas  = metas[b:b+100],
        )
    logger.success(f"Stored {len(ids)} slides in ChromaDB for [{doc_name}]")

def process_pptx(pptx_path: Path) -> dict:
    pptx_path = Path(pptx_path)
    doc_name = pptx_path.stem
    logger.info(f"Processing PPTX: {pptx_path.name}")
    cache_path = EMBEDDINGS_DIR / "text" / f"{doc_name}_ppt.json"
    if cache_path.exists():
        logger.info(f"  Skipping {pptx_path.name} — already processed")
        with open(cache_path) as f:
            data = json.load(f)
        return {
            "node_id": f"ppt_{doc_name}",
            "node_type": "PPT",
            "name": pptx_path.name,
            "path": str(pptx_path),
            "topics": data.get("topics", []),
            "embedding": data.get("embedding", []),
            "num_slides": data.get("num_slides", 0),
            "images_saved": 0,
            "slides": [],
        }

    slides = extract_slide_texts(pptx_path)
    if not slides:
        logger.warning(f"  No text found in {pptx_path.name}")
        return {"node_id": f"ppt_{doc_name}", "node_type": "PPT",
                "name": pptx_path.name, "num_slides": 0,
                "images_saved": 0, "topics": []}
    logger.info(f"  Slides with text: {len(slides)}")
    embedder   = get_embedder()
    texts      = [s["full_text"] for s in slides]
    embeddings = []
    batch_size = 16
    for i in tqdm(range(0, len(texts), batch_size),
                  desc=f"Embedding {doc_name}", leave=False):
        batch = texts[i:i+batch_size]
        vecs  = embedder.encode(batch, batch_size=batch_size,
                                show_progress_bar=False)
        embeddings.extend([v.tolist() for v in vecs])
    topics_per_slide = [extract_topics(s["full_text"]) for s in slides]
    store_in_chromadb(doc_name, pptx_path, slides, embeddings, topics_per_slide)
    all_text   = " ".join(texts)
    doc_topics = extract_topics(all_text)
    snippet    = " ".join(all_text.split()[:512])
    doc_embed  = embedder.encode(snippet).tolist()
    cache_path = EMBEDDINGS_DIR / "text" / f"{doc_name}_ppt.json"
    with open(cache_path, "w") as f:
        json.dump({
            "node_id"   : f"ppt_{doc_name}",
            "embedding" : doc_embed,
            "topics"    : doc_topics,
            "num_slides": len(slides),
        }, f)
    logger.success(f"Saved embedding cache → {cache_path}")

    saved_images = extract_images_from_pptx(pptx_path, doc_name)
    logger.info(f"  Images extracted: {len(saved_images)}")

    return {
        "node_id"     : f"ppt_{doc_name}",
        "node_type"   : "PPT",
        "name"        : pptx_path.name,
        "path"        : str(pptx_path),
        "topics"      : doc_topics,
        "embedding"   : doc_embed,
        "num_slides"  : len(slides),
        "images_saved": len(saved_images),
        "slides"      : slides,
    }

def ingest_all_pptx() -> list[dict]:
    pptx_files = list(PPT_DIR.rglob("*.pptx"))
    if not pptx_files:
        logger.warning(f"No PPTX files found in {PPT_DIR}")
        return []
    logger.info(f"Found {len(pptx_files)} PPTX files")
    results = []
    for pptx_path in tqdm(pptx_files, desc="Ingesting PPTXs"):
        try:
            results.append(process_pptx(pptx_path))
        except Exception as e:
            logger.error(f"Failed to process {pptx_path.name}: {e}")
    logger.success(f"PPTX ingestion complete: {len(results)}/{len(pptx_files)} successful")
    return results

if __name__ == "__main__":
    results = ingest_all_pptx()
    print("\n PPTX Ingestion Summary ")
    for r in results:
        print(f"  {r['name']}: {r['num_slides']} slides, {r['images_saved']} images")
    print(f"\nTotal PPTXs processed: {len(results)}")